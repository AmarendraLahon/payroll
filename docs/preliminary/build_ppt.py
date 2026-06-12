# Build AAU ERP / Payroll HLD presentation
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

DIAG = r"D:\Projects\pr\docs\diagrams"
OUT  = r"D:\Projects\pr\docs\AAU_ERP_HLD.pptx"

# palette
BLUE   = RGBColor(0x1A, 0x73, 0xE8)
DARK   = RGBColor(0x20, 0x2124)  if False else RGBColor(0x20, 0x21, 0x24)
GREY   = RGBColor(0x5F, 0x63, 0x68)
LIGHT  = RGBColor(0xF1, 0xF3, 0xF4)
GREEN  = RGBColor(0x13, 0x73, 0x33)
AMBER  = RGBColor(0xE3, 0x74, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb

def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(1.15), BLUE)
    rect(s, 0, Inches(1.15), SW, Pt(4), GREEN)
    txt(s, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8), title,
        size=26, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(s, Inches(0.52), Inches(0.74), Inches(12), Inches(0.35), kicker,
            size=12, color=RGBColor(0xD2,0xE3,0xFC))

def footer(s, n):
    txt(s, Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
        "Assam Agricultural University — ERP / Payroll  ·  High-Level Design", size=9, color=GREY)
    txt(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35), str(n), size=9,
        color=GREY, align=PP_ALIGN.RIGHT)

def bullets(s, x, y, w, h, items, size=16, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if not isinstance(it, tuple):
            lvl, t, bold, col = 0, it, False, DARK
        elif isinstance(it[0], int):           # (lvl, text)
            lvl, t = it[0], it[1]
            bold = it[2] if len(it) > 2 else False
            col  = it[3] if len(it) > 3 else DARK
        else:                                   # (text, lvl, bold[, color])
            t, lvl = it[0], it[1]
            bold = it[2] if len(it) > 2 else False
            col  = it[3] if len(it) > 3 else DARK
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl; p.space_after = Pt(gap)
        mark = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = mark + t
        r.font.size = Pt(size - lvl*2); r.font.bold = bold; r.font.color.rgb = col
        r.font.name = "Segoe UI"
    return tb

def add_image_fit(s, path, top, max_w, max_h, center_x=True, left=None):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = int(w / ar)
    if h > max_h:
        h = max_h; w = int(h * ar)
    x = int((SW - w) / 2) if center_x else left
    s.shapes.add_picture(path, x, top, width=w, height=h)
    return w, h

# ---------------- Slide 1 : Title ----------------
s = slide()
rect(s, 0, 0, SW, SH, BLUE)
rect(s, 0, Inches(4.55), SW, Pt(5), GREEN)
txt(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.2),
    "Assam Agricultural University", size=24, color=RGBColor(0xD2,0xE3,0xFC), bold=True)
txt(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.6),
    "Unified ERP & Payroll System", size=46, color=WHITE, bold=True)
txt(s, Inches(0.82), Inches(4.7), Inches(11.7), Inches(0.7),
    "High-Level Design & Requirement Analysis", size=22, color=WHITE)
txt(s, Inches(0.82), Inches(5.5), Inches(11.7), Inches(0.5),
    "Discovery Findings  ·  Proposed Architecture  ·  Phasing", size=15,
    color=RGBColor(0xD2,0xE3,0xFC))
txt(s, Inches(0.82), Inches(6.6), Inches(11.7), Inches(0.5),
    "Initial Study Report  ·  June 2026", size=13, color=RGBColor(0xAE,0xCB,0xFA))

# ---------------- Slide 2 : Agenda ----------------
s = slide(); header(s, "Agenda")
bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.2), [
    "Project context & objective",
    "Scope — three systems on one platform",
    "Proposed high-level architecture",
    "Salary → Bill → Disbursal workflow",
    "Master-driven (configurable) design",
    "Data, identifiers & integrations",
    "Delivery phasing — what we build first",
    "Open clarifications & risks",
    "Next steps",
], size=20, gap=10)
footer(s, 2)

# ---------------- Slide 3 : Context ----------------
s = slide(); header(s, "Project Context & Objective",
                    "Replace fragmented legacy systems with one web-based ERP")
bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.2), [
    ("Replacing today's separate systems:", 0, True),
    (1, "Existing HRMS — faculty/employee master, service & leave data"),
    (1, "Existing Payroll — salary structure, deductions, bill generation"),
    (1, "Finance processes — funds, budget, treasury (largely offline)"),
    ("Goal: one platform, single source of truth, full audit trail", 0, True, GREEN),
    (1, "HRMS is the only module 'properly used' today — strongest data"),
    (1, "PAN + email act as the link between sub-systems & external APIs"),
    (1, "Master-driven: funds, heads, payscales, formulas all configurable"),
], size=17, gap=8)
footer(s, 3)

# ---------------- Slide 4 : Scope ----------------
s = slide(); header(s, "Scope — Three Systems, One Platform")
cols = [
    ("Payroll  (core)", GREEN, ["Master-driven salary calc","Dynamic deductions","Payslip + annual report","Bill generation + PDF","Loans, increments, arrears"]),
    ("EFMS  (workflow)", BLUE, ["E-file movement","6 user types · 3 access levels","Full lifecycle log","Reporting-officer alerts","Access-based routing"]),
    ("FOC  (funds)", AMBER, ["Multiple fund sources","Fund → head mapping","Ceiling / allocation","Purpose-bound spend","Budget vs utilization"]),
]
cw = Inches(3.9); gap = Inches(0.25); x0 = Inches(0.7); top = Inches(1.7)
for i,(title,clr,items) in enumerate(cols):
    x = x0 + i*(cw+gap)
    rect(s, x, top, cw, Inches(0.7), clr)
    txt(s, x, top, cw, Inches(0.7), title, size=17, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x, top+Inches(0.7), cw, Inches(3.7), LIGHT)
    bullets(s, x+Inches(0.2), top+Inches(0.9), cw-Inches(0.4), Inches(3.4),
            items, size=14, gap=8)
txt(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.7),
    "+ Supporting modules (phased): Pension · GPF · Claims · Asset/Estate · Store & Purchase · Budget Mgmt",
    size=13, color=GREY)
footer(s, 4)

# ---------------- Slide 5 : HLD diagram ----------------
s = slide(); header(s, "Proposed High-Level Architecture",
                    "Layered: Users → Access → Modules → Masters → Data → External")
add_image_fit(s, os.path.join(DIAG,"hld_architecture.png"), Inches(1.35), Inches(12.6), Inches(5.5))
footer(s, 5)

# ---------------- Slide 6 : Layer walkthrough ----------------
s = slide(); header(s, "Architecture — Layer Walkthrough")
bullets(s, Inches(0.9), Inches(1.55), Inches(11.6), Inches(5.3), [
    ("Access & RBAC", 0, True, BLUE),
    (1, "Web portal + Employee Self-Service; identity via PAN/email; role-based access"),
    ("Application modules", 0, True, BLUE),
    (1, "Core Payroll (calc engine, formula engine, deductions, payslip)"),
    (1, "HRMS (employee master, service book, leave, attendance, transfer/promotion)"),
    (1, "Finance (FOC/fund, bill + PDF, treasury interface, loans)"),
    (1, "Lifecycle/compliance (pension & retirement, statutory/tax, claims)"),
    ("Cross-cutting services", 0, True, BLUE),
    (1, "EFMS workflow & approvals · notifications · audit log · reporting/analytics"),
    ("Master / Config + Data", 0, True, BLUE),
    (1, "Dynamic masters (DDO, designation, grade, fund type, budget head, payscale, formulas)"),
    (1, "Data: employee/service, payroll history, fund ledger, JSON profiles"),
], size=15, gap=5)
footer(s, 6)

# ---------------- Slide 7 : Flow diagram ----------------
s = slide(); header(s, "Salary → Bill → Disbursal Workflow",
                    "From attendance freeze to salary credit; retirement handoff to pension")
add_image_fit(s, os.path.join(DIAG,"salary_bill_flow.png"), Inches(2.2), Inches(12.6), Inches(4.2))
txt(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
    "Correction batches loop back to calc · returned bills re-enter at Accounts with reason · TA/DA stays manual (out of pay)",
    size=12, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 7)

# ---------------- Slide 8 : Master-driven ----------------
s = slide(); header(s, "Master-Driven (Configurable) Design",
                    "The spine: a rule engine, not hard-coded logic")
bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.2), [
    ("Everything is a dynamic master entry, mapped together to form 'settings':", 0, True),
    (1, "DDO · Designation · Grade · Fund Type · Budget Head · Discipline · Nature Type"),
    (1, "Payscale → mapped to designation; Fund → Head; Designation → Fund Type"),
    ("Formula engine", 0, True, GREEN),
    (1, "Formulas created in master, mapped to heads ('head assign')"),
    (1, "Deductions applied to all employees; breakup held in %"),
    (1, "Increments & corrections recomputed via formulas (changeable)"),
    ("Why it matters", 0, True, AMBER),
    (1, "New heads/rules without code changes — but needs versioning + audit"),
    (1, "Open: formula authoring rights, mid-year recompute policy (see risks)"),
], size=16, gap=6)
footer(s, 8)

# ---------------- Slide 9 : Data & integrations ----------------
s = slide(); header(s, "Data, Identifiers & Integrations")
bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.2), [
    ("Identifiers", 0, True, BLUE),
    (1, "PAN + email = cross-system link; external APIs key on PAN"),
    (1, "Risk: PAN must be unique & present for all (incl. contractual/project)"),
    ("Existing data assets", 0, True, BLUE),
    (1, "Employee personal/financial details (PAN, Aadhaar) for payroll"),
    (1, "Faculty profile, qualification, employment & research history in JSON"),
    (1, "8 appraisal formats; transfer history (dept ↔ faculty shared keys)"),
    ("External systems", 0, True, GREEN),
    (1, "Treasury · Bank · NPS/CRA · Income-Tax portal"),
    ("Retirement → handoff to pension dept; retirement age varies by cadre", 0, True),
], size=16, gap=6)
footer(s, 9)

# ---------------- Slide 10 : Phasing ----------------
s = slide(); header(s, "Delivery Phasing", "Quote the certain core; range the dependent scope")
# phase 1 box
rect(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(4.6), RGBColor(0xE6,0xF4,0xEA))
rect(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(0.65), GREEN)
txt(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(0.65), "Phase 1 — Build now (high confidence)",
    size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.95), Inches(2.55), Inches(5.5), Inches(3.6), [
    "Master/config data foundation",
    "Core payroll calc + formula engine",
    "Deductions (GPF/NPS/GIS/LIC/Tax)",
    "Payslip + annual report",
    "Bill generation + PDF + approval flow",
    "HRMS reuse (strongest existing data)",
], size=15, gap=10)
# phase 2 box
rect(s, Inches(6.8), Inches(1.7), Inches(5.85), Inches(4.6), RGBColor(0xFC,0xE8,0xD6))
rect(s, Inches(6.8), Inches(1.7), Inches(5.85), Inches(0.65), AMBER)
txt(s, Inches(6.8), Inches(1.7), Inches(5.85), Inches(0.65), "Phase 2 — Dependent (price as range)",
    size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(7.05), Inches(2.55), Inches(5.45), Inches(3.6), [
    "EFMS / full workflow engine",
    "FOC + fund-source reconciliation",
    "Pension & retirement settlement",
    "Loans (full lifecycle) & claims",
    "Offline modules (Asset/Estate, Store)",
    "External API integrations",
], size=15, gap=10)
txt(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.6),
    "Rationale: Phase 2 carries the open unknowns (fund tracking, module reality) — don't fixed-price what's untracked.",
    size=12, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 10)

# ---------------- Slide 11 : Risks / clarifications ----------------
s = slide(); header(s, "Open Clarifications & Risks", "To resolve in detailed requirement visit")
bullets(s, Inches(0.9), Inches(1.55), Inches(11.6), Inches(5.3), [
    ("Top contradictions to resolve", 0, True, AMBER),
    (1, "Fund mapping claimed vs admitted 'gap in tracking actual fund source'"),
    (1, "Salary model: pay-grade-wise vs designation-wise"),
    (1, "Modules listed vs only HRMS 'properly used' (live vs aspirational)"),
    (1, "Formula governance: who edits, versioning, mid-year recompute"),
    ("Scope boundaries to fix", 0, True, BLUE),
    (1, "Tax: TDS-only vs full compliance (Form 16/24Q)"),
    (1, "Loan: full lifecycle vs deduction line; Pension: integrated vs handoff"),
    (1, "EFMS vs bill-workflow: one engine or two; TA/DA stays out of pay"),
    ("Key risks", 0, True, RGBColor(0xC5,0x22,0x1F)),
    (1, "Legacy data quality/history span · PAN completeness · module data reality"),
], size=15, gap=5)
footer(s, 11)

# ---------------- Slide 12 : Next steps ----------------
s = slide(); header(s, "Next Steps")
bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(4.8), [
    ("Detailed requirement visit", 0, True, BLUE),
    (1, "Walk through a real monthly payroll run end-to-end"),
    (1, "Collect real filled samples: payslip, bill, formulas, fund trail"),
    (1, "Confirm module usage status (live / partial / unused / offline)"),
    ("Decisions to capture", 0, True, GREEN),
    (1, "Resolve the four contradictions; lock Phase-1 scope & sign-off"),
    (1, "Enumerate fund sources + actual disbursal path"),
    ("Output", 0, True, AMBER),
    (1, "Baselined BRD + finalized phasing → firm-up the budget estimate"),
], size=17, gap=8)
footer(s, 12)

# ---------------- Slide 13 : Thank you ----------------
s = slide()
rect(s, 0, 0, SW, SH, BLUE)
rect(s, 0, Inches(4.0), SW, Pt(5), GREEN)
txt(s, 0, Inches(2.6), SW, Inches(1.2), "Thank You", size=54, color=WHITE, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, 0, Inches(4.3), SW, Inches(0.7), "Questions & Discussion", size=22,
    color=RGBColor(0xD2,0xE3,0xFC), align=PP_ALIGN.CENTER)

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides._sldIdLst))
